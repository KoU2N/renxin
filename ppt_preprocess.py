# -*- coding: utf-8 -*-
"""[预处理模块]
负责人:lz
功能:提取ppt中的文本生成脚本和语音
"""
import os
from pptx import Presentation
from PyQt5.QtCore import QThread, pyqtSignal


class PreprocessThread(QThread):
    # ppt_path = "D:/CourseProject/Ren-robot/test.pptx"
    # ps_save_path = "D:/CourseProject/Ren-robot/lz_test.ps"
    # wav_save_path = "D:/CourseProject/Ren-robot/wavs/"
    # ppt_path = "/Users/kou2n/PycharmProjects/renxin/play/test.pptx"
    # ps_save_path = "./play/controll.ps"
    # wav_save_path = "./play/"
    # symbols = ["。", "？", "！", "；", "……", "…", ".", ";", "?", ";", "!"]
    # symbols = ["。"]
    complete_preprocess = pyqtSignal()
    update_log_view = pyqtSignal(str)

    def __init__(self):
        super().__init__()

        self.ps_save_path = "./play/controll.ps"
        self.wav_save_path = "./play/"
        self.symbols = ["。", "？", "！", "；", "……", "…", ".", ";", "?", ";", "!"]
        self.path = ''

        if os.path.exists(self.wav_save_path) is False:
            os.mkdir(self.wav_save_path)

    def run(self):
        self.remarks_list = self.obtainPPTNote(self.path)
        self.create_script(self.remarks_list)

        self.complete_preprocess.emit()

    def obtainPPTNote(self, pptPath):
        """
        :param pptPath:PPT路径
        :return:返回PPT所有备注的list列表
        """
        if pptPath.strip() == '':
            print('need a ppt \' path')
            return
        try:
            listForNote = []

            prs = Presentation(pptPath)
            for slide in prs.slides:
                # print(slide.notes_slide.notes_text_frame.text)
                listForNote.append(slide.notes_slide.notes_text_frame.text)
            # print(listForNote)
        except:
            print("opened is error for ppt")
        finally:

            # print(listForNote)
            return listForNote

    def split_sentence_by_symbol(self, sentence):
        """
        :param sentence:包含。！等终止标点的长句子
        :return: 返回分割后的字句子list列表
        """
        split_sentence = []
        pos_list = []

        for pos, char in enumerate(sentence):
            if char in self.symbols:
                pos_list.append(pos)

        temp_pos = 0
        for pos in pos_list:
            # print(sentence[temp_pos:pos+1])
            split_sentence.append(sentence[temp_pos:pos + 1])
            temp_pos = pos + 1
        # print(pos_list)
        return split_sentence

    def save_to_scripts(self, savetext, savefile):
        """
        :param savetext: 保存的内容文本字符串
        :param savefile: 保存的文件路径
        :return:
        """
        print(savetext)
        print(savefile)
        file_object = open(savefile, 'w', encoding="utf8")
        file_object.write(savetext)
        file_object.close()

    def crate_wav(self, sentence, save_wavname):
        """
        :param sentence:用于生成的语句
        :param save_wavname: 用于保存的wav文件
        :return:
        """
        from gtts import gTTS
        tts = gTTS(text=sentence, lang='ja') #lang='zh-cn'
        tts.save(self.wav_save_path + save_wavname)

    def create_script(self, remarks_list):
        """
        :param remarks_list:输入备注list
        :return:
        """
        save_text = ""
        for index, sentence in enumerate(remarks_list):
            split_sentence_list = self.split_sentence_by_symbol(sentence)
            # print(split_sentence_list)
            for sub_index, split_sentence in enumerate(split_sentence_list):
                # split_sentences = split_sentence
                # wav_name = str(index + 1) + "_" + str(sub_index + 1) + ".wav"
                wav_name = str(index + 1) + "_" + str(sub_index + 1)
                # split_wavs += wav_name+" "
                # split_wavs += wav_name
                # TODO 语句生成函数添加
                self.crate_wav(split_sentence, wav_name + ".mp3")
                print(wav_name + " | " + split_sentence + "\n")
                self.update_log_view.emit(wav_name + " | " + split_sentence + "\n")
                save_text += wav_name + " | " + split_sentence + "\n"
            # print(str(index+1)+" | "+split_sentences+" | "+split_wavs)
            # save_text += str(index+1)+" | "+split_sentences+" | "+split_wavs + "\n"
            # print(split_wavs +" | "+ split_sentences + "\n")
            # save_text += split_wavs +" | "+ split_sentences + "\n"
        self.save_to_scripts(save_text, self.ps_save_path)


# if __name__=='__main__':
#     init()
#     remarks_list = obtainPPTNote(ppt_path)
#     if remarks_list:
#         create_script(remarks_list)
#     print(sentence_list)
