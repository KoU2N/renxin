from pptx import Presentation

# load a presentation
prs = Presentation("./play/test.pptx")
# for slide in prs.slides:
#   for shapes in slide.shapes:
#     print( shapes.shape_type )
#     print( '----------------' )
#     if shapes.has_text_frame:
#       print( shapes.text )
listForNote = []
for slide in prs.slides:
    listForNote.append(slide.notes_slide.notes_text_frame.text)
print(listForNote)
